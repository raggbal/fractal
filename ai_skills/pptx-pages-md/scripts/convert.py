#!/usr/bin/env python
"""
pptx -> per-slide markdown + slide images.

For each slide in a .pptx, produce:
  - images/slide-NN.png       (rendered image of the slide)
And a single <stem>.md combining every slide (image + texts + notes).

Pipeline:
  1. python-pptx          -> extract texts and notes per slide
  2. soffice (LibreOffice) -> pptx to pdf
  3. pdf2image (poppler)   -> pdf to per-page PNG
"""

from __future__ import annotations

import argparse
import shutil
import subprocess
import sys
import tempfile
from pathlib import Path

try:
    from pptx import Presentation
except ImportError:
    print("Error: python-pptx is required. pip install -r requirements.txt", file=sys.stderr)
    sys.exit(1)

try:
    from pdf2image import convert_from_path
except ImportError:
    print("Error: pdf2image is required. pip install -r requirements.txt", file=sys.stderr)
    sys.exit(1)


def find_soffice() -> str | None:
    for cmd in ("soffice", "libreoffice"):
        path = shutil.which(cmd)
        if path:
            return path
    mac_path = "/Applications/LibreOffice.app/Contents/MacOS/soffice"
    if Path(mac_path).exists():
        return mac_path
    return None


def pptx_to_pdf(src: Path, out_dir: Path) -> Path:
    soffice = find_soffice()
    if not soffice:
        print(
            "Error: LibreOffice (soffice) not found.\n"
            "  macOS:  brew install --cask libreoffice\n"
            "  Linux:  apt install libreoffice  (or equivalent)",
            file=sys.stderr,
        )
        sys.exit(2)
    cmd = [soffice, "--headless", "--convert-to", "pdf", "--outdir", str(out_dir), str(src)]
    subprocess.run(cmd, check=True, capture_output=True)
    pdf = out_dir / f"{src.stem}.pdf"
    if not pdf.exists():
        print(f"Error: PDF not produced: {pdf}", file=sys.stderr)
        sys.exit(3)
    return pdf


def render_pdf_pages(pdf: Path, out_dir: Path, dpi: int) -> list[Path]:
    out_dir.mkdir(parents=True, exist_ok=True)
    images = convert_from_path(str(pdf), dpi=dpi)
    paths = []
    for i, img in enumerate(images, start=1):
        p = out_dir / f"slide-{i:02d}.png"
        img.save(p, "PNG")
        paths.append(p)
    return paths


def shape_iter(shapes):
    for shape in shapes:
        if shape.shape_type == 6:  # GROUP
            yield from shape_iter(shape.shapes)
        else:
            yield shape


def extract_slide_texts(slide) -> list[str]:
    texts: list[str] = []
    for shape in shape_iter(slide.shapes):
        if not shape.has_text_frame:
            continue
        for para in shape.text_frame.paragraphs:
            line = "".join(run.text for run in para.runs).strip()
            if line:
                texts.append(line)
    return texts


def extract_slide_notes(slide) -> str:
    if not slide.has_notes_slide:
        return ""
    nf = slide.notes_slide.notes_text_frame
    if not nf:
        return ""
    return nf.text.strip()


def write_combined_md(
    md_path: Path,
    src_name: str,
    slides_data: list[tuple[int, str, list[str], str]],
    slides_dir_name: str,
) -> None:
    lines = [f"# {src_name}", "", f"{len(slides_data)} slides.", ""]
    for idx, image_rel, texts, notes in slides_data:
        lines.append(f"## Slide {idx}")
        lines.append("")
        lines.append(f"![slide-{idx:02d}]({slides_dir_name}/{image_rel})")
        lines.append("")
        lines.append("### Texts")
        lines.append("")
        if texts:
            lines.extend(f"- {t}" for t in texts)
        else:
            lines.append("_(no text)_")
        lines.append("")
        lines.append("### Notes")
        lines.append("")
        lines.append(notes if notes else "_(no notes)_")
        lines.append("")
    md_path.write_text("\n".join(lines), encoding="utf-8")


def convert(src: Path, out_dir: Path, dpi: int) -> Path:
    if src.suffix.lower() != ".pptx":
        print(f"Error: not a .pptx file: {src}", file=sys.stderr)
        sys.exit(1)

    out_dir.mkdir(parents=True, exist_ok=True)
    images_dir = out_dir / "images"
    images_dir.mkdir(parents=True, exist_ok=True)

    print(f"Reading: {src.name}")
    prs = Presentation(str(src))
    slides = list(prs.slides)
    print(f"Slides: {len(slides)}")

    with tempfile.TemporaryDirectory() as tmp:
        tmp_dir = Path(tmp)
        print("Converting pptx -> pdf (LibreOffice)...")
        pdf = pptx_to_pdf(src, tmp_dir)
        print(f"Rendering pdf -> png @ {dpi} dpi...")
        images = render_pdf_pages(pdf, images_dir, dpi)

    if len(images) != len(slides):
        print(
            f"Warn: rendered {len(images)} images but pptx has {len(slides)} slides. "
            "MD/PNG pairing follows slide order.",
            file=sys.stderr,
        )

    n = max(len(images), len(slides))
    slides_data: list[tuple[int, str, list[str], str]] = []
    for i in range(1, n + 1):
        slide = slides[i - 1] if i - 1 < len(slides) else None
        texts = extract_slide_texts(slide) if slide else []
        notes = extract_slide_notes(slide) if slide else ""
        slides_data.append((i, f"slide-{i:02d}.png", texts, notes))

    md_path = out_dir / f"{src.stem}.md"
    write_combined_md(md_path, src.name, slides_data, "images")
    print(f"Output: {md_path}")
    print(f"Images dir: {images_dir}")
    return md_path


def main():
    p = argparse.ArgumentParser(description="pptx -> per-slide images + texts + notes")
    p.add_argument("source", help=".pptx file path")
    p.add_argument("-o", "--output-dir", help="Output directory (default: <source-dir>/<stem>/)")
    p.add_argument("--dpi", type=int, default=150, help="Render DPI (default: 150)")
    args = p.parse_args()

    src = Path(args.source).resolve()
    if not src.exists():
        print(f"Error: file not found: {src}", file=sys.stderr)
        sys.exit(1)

    if args.output_dir:
        out_dir = Path(args.output_dir).resolve()
    else:
        out_dir = src.parent / src.stem

    convert(src, out_dir, args.dpi)


if __name__ == "__main__":
    main()
